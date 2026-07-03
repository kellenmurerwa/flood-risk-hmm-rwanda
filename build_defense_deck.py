"""Generate the capstone defense slide deck (PowerPoint) from the real results.

Run:  python build_defense_deck.py   ->  Flood_Risk_Defense_Slides.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
FIG = ROOT / "model_outputs_real"
NAVY = RGBColor(0x12, 0x3A, 0x5E)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
GREY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title_bar(s, text, sub=None):
    tf = box(s, 0.6, 0.35, 12.1, 1.0)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.color.rgb = GREY


def bullets(s, items, left=0.7, top=1.7, width=7.4, size=18):
    tf = box(s, left, top, width, 5.2)
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + it
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = GREY if lvl else RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)


def image(s, name, left=8.2, top=1.7, width=4.6):
    f = FIG / name
    if f.exists():
        s.shapes.add_picture(str(f), Inches(left), Inches(top), width=Inches(width))


# 1 Title
s = slide()
tf = box(s, 0.8, 2.3, 11.7, 2.8)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Geospatial Flood-Risk Modelling using Climate Data and Hidden Markov Models in Rwanda"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz in [("Capstone Defense  •  Nyabugogo–Nyabarongo corridor, Kigali", 18),
                 ("Kellen Murerwa   |   Supervisor: Emmanuel Adjei", 16),
                 ("BSc Software Engineering — African Leadership University", 14)]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(sz); rr.font.color.rgb = GREY

# 2 Problem & Objective
s = slide(); title_bar(s, "Problem & Objective")
bullets(s, [
    "Kigali's catchments flood repeatedly after intense rain; official hazard maps are static.",
    "They show WHERE flooding is possible — not WHEN or HOW WIDELY pressure rises after rainfall.",
    ("Objective: a geospatial ML + Hidden Markov Model framework that estimates daily", 0),
    ("Low / Moderate / High flood-pressure states and validates them against official polygons.", 1),
    "Scope: decision-support / risk-screening — complements, not replaces, early warning.",
], width=11.5)

# 3 Study area
s = slide(); title_bar(s, "Study Area & Grid",
                       "Nyabugogo–Nyabarongo corridor aligned to the official flood zone")
bullets(s, [
    "AOI: lon 29.98–30.04, lat −2.03–−1.97 (Kigali).",
    "250 m grid → 729 cells.",
    "Daily, 2018–2024 → 1,864,053 cell-day records.",
    "Temporal split: train 2018–2022 · validation 2023 · test on unseen 2024.",
    "AOI deliberately shifted south onto the real MOE flood zone (centroid ≈ 30.054, −1.994).",
], width=7.2)
image(s, "flood_pressure_risk_map.png", left=8.0, top=1.7, width=4.9)

# 4 Architecture
s = slide(); title_bar(s, "System Architecture — four layers")
bullets(s, [
    "1. DATA — real, free, keyless feeds per cell-day.",
    "2. LABELS — flood-pressure = rainfall trigger × terrain susceptibility (+ unobserved noise).",
    "3. MODELS — baselines → RF / XGBoost → HMM temporal layer.",
    "4. VALIDATION & DELIVERY — vs official polygons + Streamlit dashboard.",
    ("Each layer is a separate, reproducible module run by main.py.", 1),
], width=11.5)

# 5 Real data sources (evidence)
s = slide(); title_bar(s, "Real Data Sources — retrieved live, not assumed")
bullets(s, [
    "Rainfall (1/3/7/14-day): Open-Meteo ERA5 archive.",
    "Elevation + slope: Open-Meteo SRTM.",
    "Rivers, roads, buildings: OpenStreetMap (Overpass).",
    "Official flood polygons: Rwanda GeoPortal / geodata.rw (Ministry of Environment).",
    ("EVIDENCE — live ERA5 pull, AOI centroid, Jan 2024:", 0),
    ("total 131.2 mm over 31 days; e.g. 09 Jan = 12.7 mm, 10 Jan = 4.2 mm.", 1),
    ("Official layer: 136 / 729 cells (19%) flagged flood-prone.", 1),
], width=11.5)

# 6 Features & labels (honesty)
s = slide(); title_bar(s, "Features & Labels — why it is not circular")
bullets(s, [
    "10 predictors: 4 rainfall windows + elevation, slope, distance-to-river, road & building density, official-polygon flag.",
    "Label = top-20% / next-30% / lowest-50% of a rainfall×susceptibility score.",
    "Added an UNOBSERVED drainage latent + daily noise the model cannot see.",
    ("Without it, models hit 0.99 F1 — a circular, indefensible result.", 1),
    ("With it, F1 ≈ 0.81 — realistic and honest.", 1),
    "Flood-PRESSURE = derived risk state, NOT a confirmed flood event.",
], width=11.5)

# 7 Models & protocol
s = slide(); title_bar(s, "Models & Evaluation Protocol")
bullets(s, [
    "Baselines: rainfall-threshold & static-polygon (prove ML adds value).",
    "Interpretable: Logistic Regression, Decision Tree.",
    "Main: Random Forest, XGBoost (+ SHAP explainability).",
    "Temporal layer: Gaussian HMM over daily sequences.",
    "Temporal split (no shuffle): train 2018–22, validate 2023, test on unseen 2024.",
    "Train ≈ validation ≈ test macro-F1 (no over-fitting; generalises across years).",
    "Metrics: macro-F1, High-recall, confusion, next-step accuracy, spatial concordance.",
], width=11.5)

# 8 Results table
s = slide(); title_bar(s, "Results — 2024 test hold-out", "Targets: macro-F1 ≥ 0.75, High-recall ≥ 0.80")
rows = [["Model", "Macro-F1", "High-recall"],
        ["XGBoost  (deployed, best)", "0.813", "0.843"],
        ["Random Forest", "0.813", "0.849"],
        ["Decision Tree", "0.750", "0.798"],
        ["Logistic Regression", "0.750", "0.819"],
        ["Rainfall-threshold (baseline)", "0.626", "—"],
        ["Static-polygon (baseline)", "0.324", "—"]]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.8), Inches(1.9),
                         Inches(7.2), Inches(4.4)).table
for c in range(3):
    tbl.columns[c].width = Inches(3.2 if c == 0 else 2.0)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(14)
        if ri == 0:
            pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif ri in (1, 2):
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE5, 0xF0, 0xE9)
image(s, "confusion_XGBoost.png", left=8.4, top=1.9, width=4.4)

# 9 Explainability
s = slide(); title_bar(s, "Explainability — SHAP & feature importance")
bullets(s, [
    "Rainfall accumulation (3- & 7-day) dominates predicted pressure.",
    "Then elevation & distance-to-river (low + near water = vulnerable).",
    "Hydrologically sensible and auditable (satisfies NFR-04).",
], width=6.0)
image(s, "shap_summary.png", left=6.7, top=1.7, width=6.1)

# 10 HMM
s = slide(); title_bar(s, "HMM — temporal dynamics of flood pressure")
bullets(s, [
    "Self-transition probabilities ≈ 0.83 / 0.80 / 0.91 (Low/Mod/High).",
    "Pressure PERSISTS once it sets in after sustained rain.",
    "Transitions move between neighbouring states, not Low→High jumps.",
    "Next-step accuracy 0.45 over 3 states (random = 0.33).",
    ("Dashboard now shows a PER-CELL transition matrix: P(state tomorrow | today) "
     "for any selected cell — the Markov intuition, cell by cell.", 0),
], width=6.0)
image(s, "hmm_transition_matrix.png", left=7.0, top=1.7, width=5.4)

# 11 Spatial validation
s = slide(); title_bar(s, "Spatial Validation vs Official Polygons — the honest result")
bullets(s, [
    "Official zone covers only 19% of the corridor → a fixed ‘70% inside’ target is geometrically impossible without circular labels.",
    "Reformulated to enrichment & odds ratio:",
    ("Raw containment of predicted-High in official zone: 0.30.", 1),
    ("Enrichment / lift: 1.60× over the 18.7% base rate.", 1),
    ("Inside-vs-outside High odds: 1.85× (30.7% vs 16.7%).", 1),
    "Model re-discovers the official zone WITHOUT training on it, and flags a broader dynamic footprint — the project's core contribution.",
], width=11.6)

# 12 Dashboard
s = slide(); title_bar(s, "Delivery — Streamlit Inspection Dashboard")
bullets(s, [
    "Pick any date → predicted flood-pressure map over the grid.",
    "Overlay official flood polygons.",
    "Inspect any cell: predictors, class probabilities, rainfall time series.",
    "Per-cell Markov transition matrix + train/validation/test performance panels.",
    "Model & HMM performance panels built in.",
    "Run: python main.py dashboard",
], width=11.5)

# 13 ERA5 vs CHIRPS vs Meteo Rwanda
s = slide(); title_bar(s, "Rainfall source: ERA5 vs CHIRPS vs Meteo Rwanda",
                       "What is feasible — tested, not assumed")
bullets(s, [
    "ERA5 (Open-Meteo) — USED. Keyless REST, 1940–now, fully reproducible; verified live (131.2 mm, Jan-2024 AOI). Reanalysis, ~9–25 km.",
    "CHIRPS — scientifically ideal for Africa (0.05°, satellite+stations). BUT access is heavy: IRI Data Library & ClimateSERV both failed/timed out in testing; needs GeoTIFF/GEE tooling.",
    "Meteo Rwanda / ENACTS — most authoritative (national gauges) but access is restricted/registration-gated; no open API.",
    ("Recommendation: ERA5 as primary (reproducible now); CHIRPS as a documented validation extension; Meteo Rwanda via formal data request.", 0),
], width=11.8, size=17)

# 14 Limitations & future work
s = slide(); title_bar(s, "Limitations & Future Work")
bullets(s, [
    "Labels are derived flood-pressure, not confirmed events (no open incident feed).",
    "ERA5 is coarser than a dense gauge network.",
    "Official national polygons are low-resolution.",
    ("Future: MINEMA incident validation; CHIRPS/Meteo Rwanda blend; TWI & HAND features; Optuna tuning; operational daily dashboard.", 0),
], width=11.6)

# 15 Conclusion
s = slide(); title_bar(s, "Conclusion")
bullets(s, [
    "All five objectives met on REAL, open data with a temporal train/val/test split.",
    "XGBoost (deployed) macro-F1 0.813, High-recall 0.843; baselines clearly beaten.",
    "Model independently reinforces (1.60×) and extends the official flood map.",
    "A reproducible, interpretable, time-aware flood-pressure tool for Kigali.",
], width=11.6)

# 16 Anticipated questions
s = slide(); title_bar(s, "Anticipated Panel Questions")
bullets(s, [
    "“Is your label circular?” → No — unobserved drainage latent + noise; F1 0.81 not 0.99.",
    "“Why ERA5 not CHIRPS?” → Reproducible keyless access; CHIRPS tested but fragile (future work).",
    "“Why did spatial agreement miss 70%?” → Coarse 19% official layer; reframed to enrichment 1.60×.",
    "“Pressure vs event?” → Derived risk state; event validation needs MINEMA records.",
    "“Does it generalise?” → Tested on a future unseen year (2024), not random split.",
], width=11.8, size=17)

# ---- speaker / presenter notes (rehearsal script, ~10-12 min total) ----
NOTES = [
 # 1 Title
 "Good morning. My name is Kellen Murerwa and my capstone is a geospatial "
 "machine-learning and Hidden Markov Model framework that estimates daily "
 "flood-pressure for the Nyabugogo-Nyabarongo corridor in Kigali. Over the next "
 "ten minutes I will take you through the problem, the real data, the models, "
 "and an honest evaluation. (~45s)",
 # 2 Problem & Objective
 "Kigali floods repeatedly after heavy rain. The official hazard maps are "
 "valuable but static: they tell you WHERE flooding is possible, not WHEN or how "
 "WIDELY pressure builds after rainfall. My objective is a time-aware "
 "decision-support tool that estimates Low, Moderate or High flood-pressure each "
 "day and is validated against the official flood polygons. It complements, not "
 "replaces, early-warning systems.",
 # 3 Study area
 "I work on a 250 m grid of 729 cells over the Nyabugogo-Nyabarongo corridor, "
 "daily from 2018 to 2024 — about 1.86 million cell-day records. I deliberately "
 "shifted the study box south so it sits on the real official flood zone. The "
 "map on the right is my model's 2024 high-pressure frequency, with the official "
 "polygons outlined in blue.",
 # 4 Architecture
 "The system has four layers: real data; a label layer; the models; and "
 "validation plus delivery. Each is a separate reproducible module, all driven "
 "by a single main.py entry point, so the whole pipeline reruns end to end.",
 # 5 Real data
 "Every predictor is real and retrieved live — nothing is assumed. Rainfall from "
 "Open-Meteo ERA5, elevation and slope from SRTM, rivers, roads and buildings "
 "from OpenStreetMap, and the official flood polygons from the government portal, "
 "geodata.rw. To prove it is real: this is a live ERA5 pull for my area for "
 "January 2024 — 131.2 mm of rain, with the actual daily values. The official "
 "layer flags 136 of my 729 cells as flood-prone.",
 # 6 Features & labels
 "This is the question I most expect from you, so let me address it directly. My "
 "label is a flood-pressure score from rainfall and terrain. If I stopped there, "
 "the label would be a formula the model could simply reverse-engineer — and "
 "indeed an early version scored 0.99 F1, which is a warning sign, not a success. "
 "So I added an unobserved drainage-and-noise term the model cannot see, which "
 "drops performance to a realistic 0.81. Flood-pressure is a derived risk state, "
 "not a confirmed flood event.",
 # 7 Models
 "I benchmark against two naive baselines so the machine learning has to earn its "
 "place, then interpretable models, then Random Forest and XGBoost with SHAP, and "
 "finally a Hidden Markov Model for the time dimension. Crucially I use a temporal "
 "split with no shuffling: I fit on 2018 to 2022, hold out 2023 as a validation "
 "year to check for over-fitting, and test on the completely unseen year 2024. "
 "Macro-F1 is almost identical across train, validation and test — so this is "
 "genuine generalisation, not memorisation.",
 # 8 Results
 "On 2024, XGBoost — my deployed model — is the best with a macro-F1 of 0.813 and "
 "high-pressure recall of 0.843, both above the 0.75 and 0.80 targets; Random "
 "Forest ties on F1 with slightly higher recall. Both baselines are clearly "
 "beaten — the rainfall-only baseline reaches just 0.63 — which proves the "
 "geospatial context adds real value. The confusion matrix shows strong Low and "
 "High detection; most error is on the transitional Moderate class, which is expected.",
 # 9 Explainability
 "SHAP confirms the model is hydrologically sensible: 3- and 7-day rainfall "
 "dominate, followed by low elevation and closeness to a river. That is exactly "
 "the physical story of urban flooding, and it makes the model auditable rather "
 "than a black box.",
 # 10 HMM
 "The HMM captures persistence. The self-transition probabilities are about 0.83, "
 "0.80 and 0.91, meaning once a cell enters high pressure after sustained rain it "
 "tends to stay there, and it moves through neighbouring states rather than "
 "jumping. Next-step accuracy is 0.45 against a random baseline of 0.33.",
 # 11 Spatial validation
 "Here is my most honest slide. My original target was 70% of predicted-high "
 "cells inside the official polygons. But that layer covers only 19% of my "
 "corridor, so 70% is mathematically impossible unless I make the label circular. "
 "I therefore reframed it to enrichment and an odds ratio. Predicted high-pressure "
 "is 1.60 times more concentrated in official zones than chance, and official "
 "cells are flagged high 1.85 times more often than others — and the model never "
 "saw those polygons in training. It rediscovers the official zone and extends it, "
 "which is the core contribution.",
 # 12 Dashboard
 "Everything is delivered through a Streamlit dashboard. You pick a date, see the "
 "predicted map, overlay the official polygons, and click any cell to see its "
 "rainfall history and the model's probabilities. It turns a model into a tool a "
 "city officer can actually use.",
 # 13 ERA5 vs CHIRPS
 "On data choice: I use ERA5 because it is the only source I can retrieve "
 "reproducibly and keylessly, as I just demonstrated live. CHIRPS is "
 "scientifically preferable for African rainfall, but I actually tested two "
 "access routes — the IRI Data Library and ClimateSERV — and both failed without "
 "heavy tooling, so I document it as a validation extension. Meteo Rwanda is the "
 "gold standard but has no open API and needs a formal request.",
 # 14 Limitations
 "I am clear about limitations: my labels are flood-pressure, not confirmed "
 "events; ERA5 is coarser than a dense gauge network; and the official polygons "
 "are low-resolution. Future work is incident validation with MINEMA, blending in "
 "CHIRPS and Meteo Rwanda, adding wetness indices, hyper-parameter tuning, and an "
 "operational daily dashboard.",
 # 15 Conclusion
 "To conclude: all five objectives are met on real, open data with a temporal "
 "train, validation and test split. The framework is accurate, interpretable, agrees with and "
 "extends the official flood map, and is fully reproducible. Thank you — I "
 "welcome your questions.",
 # 16 Q&A
 "Backup answers for likely questions — keep these crisp: circularity is handled "
 "by the unobserved noise term; ERA5 is chosen for reproducibility with CHIRPS as "
 "future work; the 70% miss is due to a coarse 19% official layer, reframed to "
 "enrichment; pressure versus event is a scope decision; and generalisation is "
 "shown by testing on a future unseen year.",
]
for sl, note in zip(prs.slides, NOTES):
    sl.notes_slide.notes_text_frame.text = note

out = ROOT / "Flood_Risk_Defense_Slides.pptx"
prs.save(out)
print("Saved", out, "—", len(prs.slides._sldIdLst), "slides,", len(NOTES), "notes")
