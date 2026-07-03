"""Build the PRACTICE deck for rehearsing the demo/defense (reference only, not a
submission artifact). Speaker notes mirror the DEMO_WALKTHROUGH.md narration so
you can practise out loud, timed to the rubric.

Run:  python demo_reference/build_practice_deck.py  ->  demo_reference/Practice_Deck.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = Path(__file__).parent
FIG = HERE.parent / "model_outputs_real"          # real figures live one level up
NAVY = RGBColor(0x12, 0x3A, 0x5E)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
GREY = RGBColor(0x44, 0x44, 0x44)
GOLD = RGBColor(0xB8, 0x86, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title_bar(s, text, sub=None):
    tf = box(s, 0.6, 0.32, 12.1, 1.0)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(14); r2.font.color.rgb = GREY


def bullets(s, items, left=0.7, top=1.7, width=7.4, size=18):
    tf = box(s, left, top, width, 5.2)
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + it
        r.font.size = Pt(size - 2 * lvl)
        r.font.color.rgb = GREY if lvl else RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)


def image(s, name, left=8.2, top=1.7, width=4.6):
    f = FIG / name
    if f.exists():
        s.shapes.add_picture(str(f), Inches(left), Inches(top), width=Inches(width))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# 1 — Title
s = slide()
tf = box(s, 0.8, 2.2, 11.7, 3.0)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Kigali Flood-Pressure Dashboard — Demo & Defense Practice"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz in [("Rehearsal deck (reference only — not for submission)", 16),
                 ("Kellen Murerwa · Supervisor: Emmanuel Adjei · ALU", 15),
                 ("Target: a clear ≤5-minute demo mapped to the rubric", 14)]:
    pp = tf.add_paragraph(); rr = pp.add_run(); rr.text = line
    rr.font.size = Pt(sz); rr.font.color.rgb = GREY
notes(s, "Open with the framing sentence: this is a decision-support tool that "
         "estimates daily flood-PRESSURE — a Low/Moderate/High risk state — for "
         "every 250 m cell in the Nyabugogo corridor. Flood-pressure is a derived "
         "risk condition, NOT a confirmed flood event. Say this before anything else.")

# 2 — The 5-minute flow
s = slide(); title_bar(s, "The 5-minute flow", "Practise this order, timed")
bullets(s, [
    "0:00 Framing — what flood-pressure is, and what it isn't.",
    "0:30 Map + date slider — the risk map that CHANGES with rainfall.",
    "1:30 Metrics + polygon overlay — agrees with the official map AND extends it.",
    "2:15 Cell inspector — features → probabilities → rainfall history.",
    "3:15 Per-cell transition matrix — the Markov contribution (spend time here).",
    "4:00 Model & splits panel — honest metrics, no over-fitting.",
    "4:45 Close — reproducible, interpretable, beats baselines.",
], width=11.6)
notes(s, "If you only get 3 minutes, keep map + cell inspector + transition matrix "
         "— that's the story that is uniquely yours. Cut the model panel to one line.")

# 3 — Problem & framing
s = slide(); title_bar(s, "Problem & framing (say it in 30s)")
bullets(s, [
    "Kigali floods after intense rain; official hazard maps are STATIC.",
    "They show WHERE flooding is possible — not WHEN or HOW WIDELY pressure rises.",
    "Our tool: daily Low / Moderate / High flood-pressure per 250 m cell.",
    ("Flood-pressure = derived risk state (rainfall + terrain + exposure).", 1),
    ("NOT a confirmed flood event. Complements early warning, doesn't replace it.", 1),
    "Real, keyless data only: ERA5 rainfall, SRTM terrain, OSM exposure, MOE polygons.",
], width=11.6)
notes(s, "The static-vs-dynamic contrast is your hook. A reviewer remembers 'their "
         "map moves with the rain, the official one can't'.")

# 4 — The product (map)
s = slide(); title_bar(s, "The product — a dynamic risk map",
                       "729 cells · green Low / orange Moderate / red High")
bullets(s, [
    "Drag the date slider dry → wet: the red High zone grows along the river corridor.",
    "Hover any cell: ID, predicted class + confidence, 3-day rain, elevation.",
    "Three colour modes: predicted state · ground-truth label · 3-day rainfall.",
    "Toggle the blue official-polygon overlay to compare.",
], width=7.2)
image(s, "flood_pressure_risk_map.png", left=8.0, top=1.7, width=4.9)
notes(s, "Move the slider LIVE. The single most persuasive moment of the demo is "
         "the red zone expanding as you step onto a high-rainfall day.")

# 5 — Results & baselines
s = slide(); title_bar(s, "Results — 2024 test hold-out", "Targets: macro-F1 ≥ 0.75, High-recall ≥ 0.80")
rows = [["Model", "Macro-F1", "High-recall", "Note"],
        ["XGBoost (deployed)", "0.813", "0.843", "best & shipped"],
        ["Random Forest", "0.813", "0.849", "ties; best recall"],
        ["Logistic Regression", "0.750", "0.819", "linear baseline"],
        ["Decision Tree", "0.750", "0.798", "clears the bar"],
        ["Rainfall-threshold (base)", "0.626", "0.902", "high recall, low precision"],
        ["Static-polygon (base)", "0.324", "0.292", "static maps inadequate"]]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.8),
                         Inches(11.9), Inches(4.4)).table
widths = [3.4, 2.0, 2.2, 4.3]
for c, w in enumerate(widths):
    tbl.columns[c].width = Inches(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13)
        if ri == 0:
            pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif ri in (1, 2):
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE5, 0xF0, 0xE9)
notes(s, "Land the value-add: beats the rainfall rule by +0.19 F1 and the static "
         "map by +0.49. The rainfall baseline's 0.90 recall is a trap — its "
         "precision is poor, which macro-F1 exposes (0.626).")

# 6 — No over-fitting (the supervisor question)
s = slide(); title_bar(s, "Train / validation / test — no over-fitting",
                       "The exact thing your supervisor asked for")
bullets(s, [
    "Temporal split, NO shuffle: train 2018–2022 · validate 2023 · test 2024.",
    "Every cell appears in all three periods; only TIME is partitioned.",
    "Macro-F1 barely moves across splits:",
    ("XGBoost: 0.804 (train) / 0.805 (val) / 0.813 (test).", 1),
    ("Random Forest: 0.813 / 0.806 / 0.813.", 1),
    "Train ≈ val ≈ test ⇒ it generalises to unseen future years, not memorising.",
    "A random split would leak future weather into training — this is the honest test.",
], width=11.6)
notes(s, "This slide is your direct answer to 'how did you split, and does it "
         "over-fit?'. Say the three numbers out loud — their closeness IS the proof.")

# 7 — Explainability
s = slide(); title_bar(s, "Confusion matrix & explainability")
bullets(s, [
    "Diagonal = correct. Strong Low & High; error concentrates on the",
    ("transitional Moderate class — off by one band, not a wild miss.", 1),
    "We optimised High-recall (0.84): a missed High day costs more than a false alarm.",
    "SHAP + feature importance: rainfall (3/7-day) dominates, then elevation & river.",
    ("Hydrologically sensible → auditable, not a black box.", 1),
], width=6.2)
image(s, "confusion_XGBoost.png", left=6.9, top=1.7, width=3.0)
image(s, "shap_summary.png", left=10.0, top=1.7, width=3.0)
notes(s, "If asked 'why is Moderate weak?' — it's the in-between band on a "
         "continuous surface; its errors land on neighbours. Point at the diagonal.")

# 8 — The Markov contribution (headline)
s = slide(); title_bar(s, "★ Per-cell Markov transitions — the contribution",
                       "Not just WHAT state — what state NEXT")
bullets(s, [
    "For a SELECTED cell: P(state tomorrow | state today), from its own history.",
    "Strong diagonal = pressure PERSISTS after sustained rain.",
    "Off-diagonal Low→Moderate→High = how a cell escalates.",
    "'—' = a state that cell never entered (shown honestly, not faked).",
    "Global HMM: self-transitions ≈ 0.83 / 0.80 / 0.91; next-step acc 0.45 (chance 0.33).",
    ("Turns a snapshot into a short-horizon forecast a manager can act on.", 1),
], width=7.0)
image(s, "hmm_transition_matrix.png", left=8.2, top=1.9, width=4.6)
notes(s, "Spend the most time here — it's what your supervisor wanted and what "
         "makes the project more than a classifier. Distinguish the per-CELL matrix "
         "(one cell's history) from the global HMM (whole corridor).")

# 9 — Spatial validation
s = slide(); title_bar(s, "Spatial validation vs the official map — the honest result")
bullets(s, [
    "Official MOE zone covers only ~19% of the corridor.",
    "So a fixed '70% of High inside the zone' is geometrically impossible without circular labels.",
    "Reframed to enrichment & odds ratio:",
    ("Enrichment / lift: 1.60× over the 18.7% base rate.", 1),
    ("Inside-vs-outside High odds: 1.85×.", 1),
    "Model RE-DISCOVERS the official zone (never trained on it) AND extends the footprint.",
], width=11.6)
notes(s, "This is your most honest slide — examiners reward candor. 'Containment "
         "0.30 looks low, but it's capped by geometry; enrichment and odds are the "
         "fair measures, and both clear the bar.'")

# 10 — Rubric self-check
s = slide(); title_bar(s, "Rubric self-check — hit every box in the demo")
rows = [["Rubric item", "Where you show it", "Number to say"],
        ["Testing Results (5)", "Model panel · testing_results/", "0.813 F1 · 10/10 tests"],
        ["Analysis (2)", "Map + leaderboard", "+0.19 vs rainfall, +0.49 vs static"],
        ["Discussion", "Persistence + spatial", "1.85× odds, agrees & extends"],
        ["Recommendations", "Limitations / future work", "incident validation, CHIRPS"],
        ["Deployment (3)", "Running app + Dockerfile", "HTTP 200 · p95 9.6 ms"],
        ["Demo / live app", "This walkthrough + URL", "record ≤5 min"]]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.8),
                         Inches(11.9), Inches(4.6)).table
for c, w in enumerate([3.3, 4.6, 4.0]):
    tbl.columns[c].width = Inches(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        pr = cell.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13)
        if ri == 0:
            pr.runs[0].font.bold = True; pr.runs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
notes(s, "Keep this as your mental checklist. If you name the number for each row "
         "during the demo, you've touched every rubric criterion explicitly.")

# 11 — Q&A prep
s = slide(); title_bar(s, "Anticipated questions — crisp answers")
bullets(s, [
    "‘Real floods?’ → No, derived pressure; incident validation is future work.",
    "‘Circular label?’ → Hidden drainage + noise; F1 0.81 not 0.99.",
    "‘Split / over-fit?’ → 3-way temporal; train≈val≈test (0.804/0.805/0.813).",
    "‘Why XGBoost?’ → Ties RF, calibrated probs + SHAP; RF has slightly higher recall.",
    "‘Why ERA5?’ → Reproducible & keyless; CHIRPS tested but fragile; Meteo Rwanda gated.",
    "‘Containment only 0.30?’ → 19% coverage caps it; enrichment 1.60× / odds 1.85×.",
    "‘Biggest limitation?’ → No incident register; labels are pressure, not events.",
], width=11.8, size=16)
notes(s, "Full answers are in DEMO_WALKTHROUGH.md §4. Rehearse Q2, Q3 and Q6 until "
         "they're automatic — those are the three most likely to be pressed.")

# 12 — Close
s = slide(); title_bar(s, "Close (one sentence)")
bullets(s, [
    "A reproducible, interpretable, time-aware flood-pressure tool for Kigali —",
    ("beats both naive baselines on real open data,", 1),
    ("agrees with and extends the official flood map,", 1),
    ("and runs as a live dashboard with a passing test suite.", 1),
    "Everything reruns end-to-end from one main.py. Thank you — questions welcome.",
], width=11.6)
notes(s, "End on reproducibility — it's the reviewer's favourite word and your "
         "strongest, most defensible claim.")

out = HERE / "Practice_Deck.pptx"
prs.save(out)
print(f"Saved {out} — {len(prs.slides._sldIdLst)} slides")
