"""
Build a concise 10-minute progress-presentation deck for the flood-risk capstone:
ML gap -> proposed solution -> current status -> working prototype, with figures.

Output: Flood_Risk_Progress_Slides.pptx  (saved in the repository root)
Run:    python build_progress_slides.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).parent
OUT = ROOT / "model_outputs_real"
SHOTS = ROOT / "Capstone final project" / "testing_results" / "screenshots"

NAVY = RGBColor(0x12, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
RED = RGBColor(0xDE, 0x2D, 0x26)
GREEN = RGBColor(0x2C, 0xA2, 0x5F)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
        p.space_after = Pt(6)
    return tb


def header(s, title, kicker):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Pt(4), RED)
    textbox(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.45),
            [(kicker.upper(), 13, RGBColor(0xBF, 0xD3, 0xE6), True)])
    textbox(s, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.65),
            [(title, 28, WHITE, True)])


def bullets(s, x, y, w, h, items, size=18):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        mark = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = mark + txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = color
        r.font.name = "Calibri"
        p.space_after = Pt(8)
    return tb


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def add_image_fit(s, img, x, y, w, h, caption=None):
    from PIL import Image
    img = str(img)
    iw, ih = Image.open(img).size
    box_ar = w / h; img_ar = iw / ih
    if img_ar > box_ar:
        nw = w; nh = Emu(int(w / img_ar))
    else:
        nh = h; nw = Emu(int(h * img_ar))
    px = x + Emu(int((w - nw) / 2)); py = y + Emu(int((h - nh) / 2))
    s.shapes.add_picture(img, px, py, width=nw, height=nh)
    if caption:
        textbox(s, x, y + h, w, Inches(0.35),
                [(caption, 11, GREY, False)], align=PP_ALIGN.CENTER)


# ----------------------------------------------------------------- 1. TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.35), SW, Pt(5), RED)
textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.4),
        [("Geospatial Flood-Risk Modelling using Climate Data", 36, WHITE, True),
         ("and Hidden Markov Models in Rwanda", 36, WHITE, True),
         ("Nyabugogo–Nyabarongo corridor, Kigali", 20, RGBColor(0xBF, 0xD3, 0xE6), False)])
textbox(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(1.6),
        [("Progress presentation  ·  ML gap → solution → status → prototype", 18, LIGHT, False),
         ("Kellen Murerwa   |   Supervisor: Emmanuel Adjei", 16, WHITE, True),
         ("BSc Software Engineering · African Leadership University", 14, RGBColor(0xBF, 0xD3, 0xE6), False)])
notes(s, "Good [morning]. I'm Kellen Murerwa. My capstone builds a geospatial "
         "machine-learning system that estimates daily flood-pressure for a "
         "flood-prone corridor in Kigali, using only open climate and map data. "
         "In the next 10 minutes I'll cover the ML gap, my solution, where I am "
         "now, and a working prototype. (~30s)")

# ----------------------------------------------------------------- 2. ML GAP
s = slide()
header(s, "The ML gap we address", "Problem")
bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.4), [
    ("Kigali's existing flood-risk tools are STATIC hazard maps — they cover only "
     "~19% of the study corridor and never change with the weather.", 0, NAVY),
    ("They give no time-aware, day-to-day view of where flood pressure is "
     "building after heavy rain.", 1, GREY),
    ("ML flood-susceptibility studies exist, but in this setting they rarely:", 0, NAVY),
    ("combine supervised geospatial ML with a temporal (sequence) model,", 1, GREY),
    ("rely only on open / free data for a data-scarce African urban catchment,", 1, GREY),
    ("or ship a validated, usable artefact — most stop at an accuracy table.", 1, GREY),
    ("→ The gap: a reproducible, time-aware ML system from OPEN data that beats "
     "static maps and rule baselines — and is validated against the official "
     "flood polygons.", 0, RED),
])
notes(s, "The problem: Kigali's flood tools are static hazard maps. They're "
         "authoritative but cover only ~19% of my corridor and never change — "
         "so they can't tell you where pressure is building after a storm. ML "
         "studies exist, but they rarely pair geospatial ML with a temporal "
         "model, rarely stick to open data for an African city, and rarely ship "
         "something usable. That intersection is my gap. (~1.5 min)")

# ------------------------------------------------------ 3. THE GAP, VISUALISED
s = slide()
header(s, "Static map vs. dynamic prediction", "The gap, visualised")
add_image_fit(s, OUT / "static_vs_dynamic.png",
              Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.0))
textbox(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
        [("Left: the fixed official zone (same every day). Right: on a wet day "
          "the model flags a far broader, rainfall-driven footprint — the risk "
          "the static map misses.", 14, GREY, False)], align=PP_ALIGN.CENTER)
notes(s, "Here's the gap in one picture. On the left is the official hazard map "
         "— a fixed blue zone, ~19% of the area, identical every day of the year. "
         "On the right is my model on a wet day: it still concentrates risk in "
         "and around the official zone, but it flags a much wider rainfall-driven "
         "footprint the static map simply can't show. Same place, but a "
         "time-aware view. (~1 min)")

# ----------------------------------------------------------------- 3. SOLUTION
s = slide()
header(s, "Proposed solution", "Approach")
bullets(s, Inches(0.5), Inches(1.45), Inches(7.4), Inches(5.6), [
    ("A 250 m grid over the corridor (729 cells), one row per cell-per-day, "
     "2018–2024 (1.86 M cell-days).", 0, NAVY),
    ("Fuse four open, keyless sources into 10 features:", 0, NAVY),
    ("rainfall accumulations (1/3/7/14-day) — Open-Meteo ERA5,", 1, GREY),
    ("terrain (elevation, slope) — SRTM; exposure (river, roads, buildings) — OSM,", 1, GREY),
    ("official flood-zone membership — Rwanda GeoPortal / MOE.", 1, GREY),
    ("Predict daily Low / Moderate / High flood-pressure states.", 0, NAVY),
    ("Benchmark LR/DT, Random Forest, XGBoost vs. rainfall-threshold & "
     "static-polygon baselines; add an HMM temporal layer.", 0, NAVY),
    ("Honest test: temporal split — train 2018–22, validate 2023, test on unseen "
     "2024; then spatial validation + a Streamlit dashboard.", 0, BLUE),
])
rect(s, Inches(8.2), Inches(1.6), Inches(4.6), Inches(4.9), LIGHT)
add_image_fit(s, OUT / "xgb_feature_importance.png",
              Inches(8.35), Inches(1.75), Inches(4.3), Inches(4.4),
              "Engineered features ranked by importance")
notes(s, "My solution: lay a 250 m grid over the corridor — 729 cells, one row "
         "per cell per day, seven years, 1.86 million rows. I fuse four open, "
         "keyless sources into ten features: rainfall accumulations from "
         "Open-Meteo ERA5, terrain from SRTM, exposure — rivers, roads, "
         "buildings — from OpenStreetMap, and official flood-zone membership "
         "from the Rwanda GeoPortal. I predict Low/Moderate/High pressure, "
         "benchmark Random Forest and XGBoost against rule and static-map "
         "baselines, add an HMM temporal layer, and — crucially — test on an "
         "unseen year, 2024. The chart shows rainfall features dominate, which "
         "is hydrologically sensible. (~2 min)")

# ----------------------------------------------------------------- 4. STATUS
s = slide()
header(s, "Current status", "Where we are")
done = [
    "Data pipeline COMPLETE — 1.86 M real cell-days, fully reproducible (cached, offline-friendly).",
    "Models trained & benchmarked on the 2024 test hold-out — targets MET.",
    "XGBoost (deployed) macro-F1 0.813 / High-recall 0.843;  Random Forest 0.813 / 0.849.",
    "Train ≈ val ≈ test macro-F1 — no over-fitting. Beats baselines: +0.19 F1 over rainfall, +0.49 over static map.",
    "Spatial validation DONE — enrichment 1.60× (≥1.4 ✅), inside/outside odds 1.85× (≥1.5 ✅).",
    "HMM temporal layer integrated (descriptive; supervised layer is the predictor).",
    "Working dashboard + 10/10 automated tests; live-query latency p95 ≈ 9.6 ms (budget 2 s).",
    "Clean GitHub repo + deployment artifacts (Docker / Streamlit Cloud / Render) ready.",
]
bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
        [(t, 0, GREEN if i in (0, 1, 4, 6, 7) else NAVY) for i, t in enumerate(done)],
        size=18)
notes(s, "Status: the build is essentially done. The data pipeline is complete "
         "and reproducible. Models are trained and tested on 2024 — and the "
         "proposal targets are met: XGBoost — my deployed model — hits 0.81 "
         "macro-F1 and 0.84 high-recall, beating the rainfall rule by 0.19 F1 and "
         "the static map by 0.49, and train, validation and test scores match so "
         "it isn't over-fitting. Spatial validation passes both targets — 1.6x "
         "enrichment and 1.85x odds inside the official zones. The HMM is integrated as an "
         "explanatory layer. And there's a working dashboard, a full passing "
         "test suite, sub-10-millisecond queries, and a clean repo with "
         "deployment ready. (~1.5 min)")

# --------------------------------------------------------- 5. PROTOTYPE (dashboard)
s = slide()
header(s, "Working prototype — the dashboard", "Demo")
add_image_fit(s, SHOTS / "dashboard_preview.png",
              Inches(0.5), Inches(1.45), Inches(8.6), Inches(5.4),
              "Pick a date → predicted flood-pressure map, metrics, per-cell inspector")
bullets(s, Inches(9.3), Inches(1.6), Inches(3.7), Inches(5.4), [
    ("Streamlit app", 0, NAVY),
    ("Any date → Low/Mod/High map over 729 cells", 1, GREY),
    ("Toggle official flood polygons", 1, GREY),
    ("Inspect a cell: features + class probabilities", 1, GREY),
    ("Per-cell Markov transition matrix", 1, GREY),
    ("Rainfall series, train/val/test & HMM panels", 1, GREY),
    ("Boots headless, serves HTTP 200", 0, GREEN),
    ("p95 ≈ 9.6 ms / query", 0, GREEN),
], size=15)
notes(s, "Here's the prototype — a Streamlit dashboard. You pick any date and "
         "see the predicted Low/Moderate/High map over all 729 cells. You can "
         "overlay the official polygons, click any cell to see its features and "
         "the model's class probabilities, and view rainfall trends and the HMM "
         "panel. It boots cleanly and answers a query in under 10 milliseconds, "
         "so it's genuinely interactive. (~1.5 min — demo live here if time)")

# --------------------------------------------------------- 6. PROTOTYPE (evidence)
s = slide()
header(s, "Working prototype — it works on real & varied data", "Evidence")
add_image_fit(s, SHOTS / "confusion_XGBoost.png",
              Inches(0.4), Inches(1.5), Inches(4.1), Inches(4.6),
              "XGBoost confusion (2024 hold-out)")
add_image_fit(s, SHOTS / "demo_data_values.png",
              Inches(4.6), Inches(1.5), Inches(4.5), Inches(4.6),
              "Prediction responds to different data values")
add_image_fit(s, SHOTS / "benchmark_performance.png",
              Inches(9.2), Inches(1.5), Inches(3.9), Inches(4.6),
              "Latency well inside the 2 s budget")
notes(s, "And it's tested, not just demoed. Left: the confusion matrix on the "
         "2024 hold-out — errors sit on neighbouring states, not wild misses. "
         "Middle: a functional test showing predictions shift correctly as I "
         "change the inputs — dry high ground stays Low, heavy rain on low "
         "ground goes High. Right: a performance benchmark — latency stays well "
         "under the 2-second budget even on a single thread. Ten of ten "
         "automated tests pass. (~1.5 min)")

# ----------------------------------------------------------------- 7. CLOSING
s = slide()
header(s, "Summary & next steps", "Wrap-up")
bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0), [
    ("Thesis demonstrated: open climate + geospatial data + ensemble ML give a "
     "time-aware flood-pressure view that beats static maps and rule baselines.", 0, NAVY),
    ("End-to-end and reproducible: data → models → validation → live dashboard.", 0, NAVY),
], size=19)
textbox(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.5),
        [("Next steps", 20, RED, True)])
bullets(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.2), [
    ("Ground-truth against real flood-incident records (biggest validity gain).", 0, GREY),
    ("Calibrate probabilities → a meaningful Moderate band; feed forecast rainfall "
     "for a 1–7 day outlook.", 0, GREY),
    ("Operationalise: scheduled refresh, drift monitoring, alert push; finalise hosting.", 0, GREY),
], size=17)

notes(s, "To wrap up: I've shown the thesis holds — open climate and geospatial "
         "data plus ensemble ML give a time-aware flood view that beats static "
         "maps and rule baselines, end-to-end and reproducible. Next, I'll "
         "ground-truth against real flood-incident records, calibrate the "
         "probabilities and feed forecast rainfall for a forward-looking "
         "outlook, and finalise hosting. Thank you — happy to take questions. "
         "(~1 min)")

path = ROOT / "Flood_Risk_Progress_Slides.pptx"
prs.save(str(path))
print(f"Saved {path}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
