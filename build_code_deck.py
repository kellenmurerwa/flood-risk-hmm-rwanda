"""
Build a CODE-WALKTHROUGH revision deck: real code snippets from the flood-risk
product + plain-English explanations, the tech/ML stack, and a 'what to explain'
cheat sheet. Companion to Flood_Risk_Progress_Slides.pptx.

Output: Flood_Risk_Code_Walkthrough.pptx  (saved in the repository root)
Run:    python build_code_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent

NAVY = RGBColor(0x12, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
RED = RGBColor(0xDE, 0x2D, 0x26)
GREEN = RGBColor(0x2C, 0xA2, 0x5F)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
CODEBG = RGBColor(0x1E, 0x29, 0x39)
CODEFG = RGBColor(0xE6, 0xED, 0xF3)
CODEKEY = RGBColor(0x7E, 0xC6, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(): return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
        p.space_after = Pt(5)
    return tb


def header(s, title, kicker):
    rect(s, 0, 0, SW, Inches(1.05), NAVY)
    rect(s, 0, Inches(1.05), SW, Pt(4), RED)
    textbox(s, Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.4),
            [(kicker.upper(), 12, RGBColor(0xBF, 0xD3, 0xE6), True)])
    textbox(s, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.6),
            [(title, 26, WHITE, True)])


def code_panel(s, x, y, w, h, code, fontsize=12.5):
    rect(s, x, y, w, h, CODEBG)
    tb = s.shapes.add_textbox(x + Inches(0.12), y + Inches(0.08),
                              w - Inches(0.24), h - Inches(0.16))
    tf = tb.text_frame; tf.word_wrap = False
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0); p.line_spacing = 1.0
        is_comment = line.lstrip().startswith("#")
        r = p.add_run(); r.text = line if line else " "
        r.font.name = "Consolas"; r.font.size = Pt(fontsize)
        r.font.color.rgb = RGBColor(0x8B, 0x9C, 0xB0) if is_comment else CODEFG
    return tb


def bullets(s, x, y, w, h, items, size=15):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + txt
        r.font.size = Pt(size - lvl); r.font.color.rgb = color; r.font.name = "Calibri"
        p.space_after = Pt(6)
    return tb


def code_slide(title, kicker, code, explain, fontsize=12.5):
    s = slide(); header(s, title, kicker)
    code_panel(s, Inches(0.45), Inches(1.3), Inches(7.7), Inches(5.7), code, fontsize)
    bullets(s, Inches(8.35), Inches(1.4), Inches(4.55), Inches(5.6), explain)
    return s


def notes(s, t): s.notes_slide.notes_text_frame.text = t


# ------------------------------------------------------------------ TITLE
s = slide()
rect(s, 0, 0, SW, SH, NAVY); rect(s, 0, Inches(4.2), SW, Pt(5), RED)
textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.2),
        [("Code Walkthrough & Tech Stack", 38, WHITE, True),
         ("Geospatial Flood-Risk Modelling (Climate Data + HMM), Kigali", 20,
          RGBColor(0xBF, 0xD3, 0xE6), False)])
textbox(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(1.4),
        [("Revision companion: how the product is built, function by function",
          18, LIGHT, False),
         ("Kellen Murerwa  ·  Supervisor: Emmanuel Adjei  ·  ALU BSc SE", 15, WHITE, True)])

# --------------------------------------------------- 1. ORCHESTRATOR
code_slide("1 · Pipeline orchestrator — main.py", "Build-up · CLI entry point",
"""STAGES = {
  "build":    "build_real_dataset.py",  # real data -> parquet
  "polygons": "add_official_polygons.py",# official MOE polygons
  "train":    "train_real.py",          # models + HMM (2024 hold-out)
  "evaluate": "evaluate_spatial.py",    # spatial validation + map
}

# one command runs the whole thing, reproducibly:
#   python main.py all        build -> train -> evaluate
#   python main.py dashboard  launch the Streamlit app

if args.stage == "all":
    for s in ["build", "train", "evaluate"]:
        run_module(STAGES[s])""",
[("Single CLI entry point; each stage is an independent, re-runnable module.", 0, NAVY),
 ("`all` chains build → train → evaluate; `dashboard` serves the app.", 0, GREY),
 ("API responses cached in data_cache/, so re-runs are fast & offline.", 0, GREY),
 ("This is the slide to open on — it frames the whole architecture.", 0, BLUE)])

# --------------------------------------------------- 2. DATA ACQUISITION
code_slide("2 · Data acquisition — real, free, keyless", "Functionality · sources",
'''# rainfall: Open-Meteo ERA5 reanalysis archive
requests.get("https://archive-api.open-meteo.com/v1/archive",
  params={"latitude": la, "longitude": lo,
          "start_date": "2018-01-01", "end_date": "2024-12-31",
          "daily": "precipitation_sum"})

# OpenStreetMap via Overpass (rotates mirrors + User-Agent)
requests.post(url, data={"data": query}, headers=OVERPASS_HEADERS)

# official flood zones: Rwanda GeoPortal (ArcGIS REST)
requests.get(SERVICE, params={"where": "objectid IN (1,2,3)",
                              "outSR": "4326", "f": "geojson"})''',
[("Four open sources, no API keys: ERA5 rainfall, SRTM elevation, OSM "
  "exposure, MOE flood polygons.", 0, NAVY),
 ("Overpass call hardened: User-Agent + mirror rotation + JSON validation.", 0, GREY),
 ("Honest caveat: rainfall is ERA5 reanalysis (substituted for Meteo Rwanda).", 0, RED)])

# --------------------------------------------------- 3. FEATURE ENGINEERING
code_slide("3 · Feature engineering", "Functionality · 10 features",
"""# rainfall accumulation windows (rolling sums)
"rainfall_3d_mm":  s.rolling(3,  min_periods=1).sum()
"rainfall_7d_mm":  s.rolling(7,  min_periods=1).sum()
"rainfall_14d_mm": s.rolling(14, min_periods=1).sum()

# distance to nearest river: spatial KD-tree (fast NN)
rtree = cKDTree(river_xy)
dist, _ = rtree.query(cell_xy)

# slope from the elevation grid's gradient
gy, gx = np.gradient(Z, CELL_M, CELL_M)
slope  = np.degrees(np.arctan(np.hypot(gx, gy)))

# road length density & building count density per km^2""",
[("Turns raw data into flood-relevant signal: multi-day rain build-up, "
  "terrain, proximity & exposure.", 0, NAVY),
 ("`cKDTree` (SciPy) = efficient nearest-river distance for 729 cells.", 0, GREY),
 ("Slope derived from the SRTM elevation surface via numpy gradient.", 0, GREY),
 ("SHAP later confirms the rainfall-accumulation features dominate.", 0, GREEN)])

# --------------------------------------------------- 4. LABELS
code_slide("4 · Flood-pressure labels", "Functionality · the target",
'''# rainfall trigger (saturating) x static susceptibility
trig = 0.6*np.tanh(rain3/60) + 0.4*np.tanh(rain7/110)
base = trig * (0.12 + 0.88*susc)     # rain GATED by terrain

# + UNOBSERVED drainage latent + daily noise
#   (deliberately NOT given to the model as features)
score = base + cell_latent + daily_noise

# 50 / 30 / 20  ->  Low / Moderate / High
q1, q2 = np.quantile(score, [0.50, 0.80])
state = np.where(score <= q1, "Low",
        np.where(score <= q2, "Moderate", "High"))''',
[("Label = rainfall trigger × terrain susceptibility (heavy rain on high, "
  "well-drained ground stays Low).", 0, NAVY),
 ("Latent drainage + noise are hidden from the model on purpose — keeps the "
  "task realistic and NON-circular.", 0, RED),
 ("Expect this question: 'are your labels derived?' Yes — and the validity "
  "limit is owned in the report.", 0, BLUE)])

# --------------------------------------------------- 5. MODELS + SPLIT
code_slide("5 · Models & temporal train/val/test split — train_real.py", "ML tools · training",
'''tr = df[df.date.dt.year <= 2022]   # train  2018-2022
va = df[df.date.dt.year == 2023]   # validate 2023 (over-fit check)
te = df[df.date.dt.year == 2024]   # test on UNSEEN 2024

models = {
 "LogisticRegression": make_pipeline(StandardScaler(),
                        LogisticRegression(max_iter=2000)),
 "DecisionTree": DecisionTreeClassifier(max_depth=8),
 "RandomForest": RandomForestClassifier(n_estimators=300,
                 min_samples_leaf=5, max_leaf_nodes=4096, n_jobs=2),
 "XGBoost": xgb.XGBClassifier(n_estimators=400, max_depth=5,
            learning_rate=0.07, tree_method="hist",
            objective="multi:softprob", num_class=3),
}
# + rainfall-threshold & static-polygon baselines''',
[("3-way temporal split (not random) = the honest test; 2024 is never seen in training.", 0, NAVY),
 ("Train ≈ validation ≈ test macro-F1 → no over-fitting (reported per split).", 0, GREY),
 ("scikit-learn (LR/DT/RF) + XGBoost; hyper-params are explicit, not defaults.", 0, GREY),
 ("Result: XGBoost (deployed) 0.813 / 0.843; RF 0.813 / 0.849.", 0, GREEN)])

# --------------------------------------------------- 6. HMM
code_slide("6 · HMM temporal layer", "ML tools · time dynamics",
'''from hmmlearn import hmm

# per-cell daily rainfall sequences -> 3 latent states
model = hmm.GaussianHMM(n_components=3,
            covariance_type="diag", n_iter=200)
model.fit(Xs, lengths)

decoded = model.predict(Xs, [L])   # Viterbi state per day
T = model.transmat_                # daily transition matrix
#   High->High = 0.91  (flood pressure is persistent)

# next-step accuracy: argmax(T[state_t]) vs truth_{t+1}''',
[("Models day-to-day DYNAMICS of flood pressure (not just a static snapshot).", 0, NAVY),
 ("Transition matrix shows strong persistence & neighbour-only moves.", 0, GREY),
 ("Dashboard renders a PER-CELL transition matrix — P(state tomorrow | today).", 0, GREY),
 ("Honest framing: HMM is explanatory; the supervised layer is the predictor "
  "(next-step acc 0.45).", 0, RED)])

# --------------------------------------------------- 7. EVALUATION
code_slide("7 · Evaluation & targets", "Functionality · validation",
'''macro_f1    = f1_score(yte, p, average="macro")
high_recall = recall_score(yte, p, labels=[2], average="macro")

# spatial validation vs OFFICIAL polygons (evaluate_spatial.py)
lift = P(official_zone | predicted_High) / P(official_zone)
odds = high_rate_inside_zone / high_rate_outside_zone

targets_met = {
  "macro_f1>=0.75":   True,    # 0.813
  "high_recall>=0.80": True,   # 0.843
}   # enrichment 1.60x (>=1.4), odds 1.85x (>=1.5)''',
[("Reports F1, accuracy, recall + a spatial check against the official map.", 0, NAVY),
 ("Enrichment & odds-ratio chosen because polygons cover only ~19% of cells.", 0, GREY),
 ("All proposal targets met; results written to results_summary.json.", 0, GREEN)])

# --------------------------------------------------- 8. DASHBOARD
code_slide("8 · The product — dashboard.py", "Functionality · the app",
'''import streamlit as st, pydeck as pdk

@st.cache_resource
def load_model():
    b = joblib.load("model_outputs_real/xgboost_model.joblib")
    return b["model"], b["features"]

day["pred_i"] = model.predict(day[feat].values)   # per cell
proba = model.predict_proba(day[feat].values)     # confidence

pdk.Layer("ScatterplotLayer", day,        # the live map
    get_fill_color="color", get_radius=110, pickable=True)''',
[("Streamlit app: pick a date → live Low/Mod/High map over 729 cells.", 0, NAVY),
 ("`@st.cache_resource` loads the model once; `pydeck` renders the map.", 0, GREY),
 ("Click a cell → features + class probabilities; p95 query ≈ 9.6 ms.", 0, GREEN),
 ("Demo this live, or play flood_pressure_animation.gif.", 0, BLUE)])

# --------------------------------------------------- 9. TECH STACK
s = slide(); header(s, "9 · Technologies, languages & ML tools", "Stack at a glance")
groups = [
 ("Language", "Python 3.11", BLUE),
 ("Data handling", "pandas · numpy · pyarrow (Parquet) · SciPy (cKDTree)", NAVY),
 ("Machine learning", "scikit-learn (LogReg, DecisionTree, RandomForest, "
  "metrics, scaling) · XGBoost · hmmlearn (HMM) · SHAP", GREEN),
 ("Data sources / APIs", "Open-Meteo (ERA5 rainfall, SRTM elevation) · "
  "OpenStreetMap / Overpass · Rwanda GeoPortal (ArcGIS REST)", NAVY),
 ("Visualisation / app", "matplotlib · Streamlit · pydeck", BLUE),
 ("Persistence / infra", "joblib (model artefacts) · Docker · GitHub · "
  "Streamlit Cloud / Render", NAVY),
 ("Testing / quality", "pytest (10 tests) · temporal hold-out · spatial validation", GREEN),
]
y = 1.35
for label, val, col in groups:
    rect(s, Inches(0.5), Inches(y), Inches(3.0), Inches(0.66), col)
    textbox(s, Inches(0.6), Inches(y + 0.12), Inches(2.85), Inches(0.5),
            [(label, 14, WHITE, True)])
    textbox(s, Inches(3.7), Inches(y + 0.06), Inches(9.1), Inches(0.62),
            [(val, 14, NAVY, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78

# --------------------------------------------------- 10. WHAT TO EXPLAIN
s = slide(); header(s, "10 · What to explain (defense cheat-sheet)", "Be ready for")
bullets(s, Inches(0.5), Inches(1.35), Inches(12.4), Inches(5.6), [
 ("Architecture: one pipeline, four stages (main.py) — start here.", 0, NAVY),
 ("Data provenance: the four open sources + why ERA5 instead of Meteo Rwanda.", 0, NAVY),
 ("Feature engineering: rainfall accumulation windows, distance-to-river "
  "(cKDTree), slope from elevation gradient, densities.", 0, NAVY),
 ("Labels: rainfall trigger × susceptibility + hidden latent/noise — why it's "
  "NOT circular (the most likely hard question).", 0, RED),
 ("Why a temporal split (train 2018–22, validate 2023, test 2024) rather than a random split.", 0, NAVY),
 ("Model choice: RF/XGBoost vs baselines; what the numbers mean (macro-F1, "
  "High-recall) and why recall is weighted.", 0, NAVY),
 ("Spatial validation: enrichment 1.60× & odds 1.85× — why not raw containment "
  "(polygons cover only ~19%).", 0, NAVY),
 ("HMM: what it adds (dynamics/persistence) and its honest limit vs supervised.", 0, GREY),
 ("The dashboard + tests + deployment = a usable, validated artefact, not just a table.", 0, GREEN),
])

path = ROOT / "Flood_Risk_Code_Walkthrough.pptx"
prs.save(str(path))
print(f"Saved {path}  ({len(prs.slides._sldIdLst)} slides)")
